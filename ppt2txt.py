
"""
# ppt2txt.py
# script recursively extractis text from all pptx files
# in current directory - and saves to a text file 
# or you can provide directories on command line (space-separaed)
""" 

# pip install python-pptx
import os, sys, glob
from pptx import Presentation


if len(sys.argv) <= 1:
    mypath = os.path.dirname(os.path.realpath(__file__))
    print("looking for pptx files in script directory: ", mypath)
    mydirs = [mypath]
else:
    mydirs = []
    for argpath in sys.argv[1:]:
        if not os.path.isdir(argpath):
            print(f"{argpath} is not a directory, exiting ...")
            sys.exit()
        mypath = os.path.realpath(argpath)
        print(argpath, "   =>  ", mypath)
        mydirs += [mypath]

for dir_path in mydirs:
    print(dir_path)
    basedir = dir_path.split("/")[-1]
    fname_out = dir_path + "/" + basedir+"_ppt.txt"

    ss=""
    counter = 0
    myfiles = glob.glob(dir_path+'/**/*pptx', recursive=True)
    for myfile in sorted(myfiles):
        if "~$" in myfile:
            continue
        if "OLD/" in myfile:
            continue
        counter += 1
        if counter > 3:
            pass # break
        prs = Presentation(myfile)
        print("-"*20, myfile)
        myfile_short = myfile.split(dir_path)[1][1:]
        # print(myfile_short,"\n")
        # print("----------------------")
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = shape.text
                    txt_arr = txt.strip().split("\n")
                    txt_arr = [x.strip() for x in txt_arr]
                    txt_arr = [x.replace("\r"," ") for x in txt_arr]
                    txt_arr = [myfile_short+" : " + x + "\n" for x in txt_arr]
                    txt = "".join(txt_arr)
                    ss += txt
                    print(txt)
    print("----------------------")
    print("writing to file ", fname_out)
    fh = open(fname_out, "w")
    fh.write(ss)
    fh.close()
    print("----------------------")
